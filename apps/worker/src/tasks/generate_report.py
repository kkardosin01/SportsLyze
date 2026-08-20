"""Geração de relatório (PDF, slides ou áudio) por partida ou por jogador.
Estrutura pronta para futuros formatos — ver `reports.format` no schema.
Os três formatos reaproveitam o mesmo `report_context` (stats por jogador,
clube, adversário, data), só muda o motor de renderização de cada um."""

import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from celery import shared_task
from gtts import gTTS
from jinja2 import Environment, FileSystemLoader
from pptx import Presentation
from pptx.util import Inches
from weasyprint import HTML

from src.core.config import get_settings
from src.db.session import get_supabase_admin

TEMPLATES_DIR = Path(__file__).resolve().parent.parent / "templates"
_jinja_env = Environment(loader=FileSystemLoader(str(TEMPLATES_DIR)))

_CONTENT_TYPES = {
    "pdf": "application/pdf",
    "slides": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "audio": "audio/mpeg",
}
_EXTENSIONS = {"pdf": "pdf", "slides": "pptx", "audio": "mp3"}


@shared_task(name="tasks.generate_match_report")
def generate_match_report(match_id: str, created_by: str, report_format: str = "pdf") -> str:
    if report_format not in _CONTENT_TYPES:
        raise ValueError(f"Formato de relatório não suportado: {report_format}")

    db = get_supabase_admin()
    settings = get_settings()

    match = db.table("matches").select("*, clubs(name)").eq("id", match_id).single().execute().data
    videos = db.table("videos").select("id").eq("match_id", match_id).execute().data or []
    video_ids = [v["id"] for v in videos]

    stats_rows = (
        db.table("player_match_stats")
        .select("*, detected_players(label, athlete_id, athletes(full_name))")
        .in_("video_id", video_ids)
        .execute()
        .data
        if video_ids
        else []
    )

    players = [
        {
            "name": (row["detected_players"].get("athletes") or {}).get("full_name")
            or row["detected_players"]["label"],
            "distance_km": row["distance_km"] or 0,
            "avg_speed_kmh": row["avg_speed_kmh"] or 0,
            "max_speed_kmh": row["max_speed_kmh"] or 0,
        }
        for row in stats_rows
    ]

    context = {
        "club_name": (match.get("clubs") or {}).get("name", "Atleta independente"),
        "opponent_name": match.get("opponent_name"),
        "match_date": match.get("match_date") or "",
        "generated_at": datetime.now(timezone.utc).strftime("%d/%m/%Y %H:%M"),
        "players": players,
    }

    with tempfile.TemporaryDirectory(prefix="sportslyze_report_") as tmp_dir:
        if report_format == "pdf":
            file_bytes = _render_pdf(context, tmp_dir)
        elif report_format == "slides":
            file_bytes = _render_slides(context, tmp_dir)
        else:
            file_bytes = _render_audio(context, tmp_dir)

        extension = _EXTENSIONS[report_format]
        storage_path = f"matches/{match_id}/relatorio_{int(datetime.now(timezone.utc).timestamp())}.{extension}"
        db.storage.from_(settings.supabase_storage_bucket_reports).upload(
            storage_path, file_bytes, {"content-type": _CONTENT_TYPES[report_format]}
        )

    report = (
        db.table("reports")
        .insert(
            {
                "scope": "match",
                "match_id": match_id,
                "format": report_format,
                "storage_path": storage_path,
                "created_by": created_by,
            }
        )
        .execute()
        .data[0]
    )

    return report["id"]


def _render_pdf(context: dict[str, Any], tmp_dir: str) -> bytes:
    template = _jinja_env.get_template("report_match.html")
    html_content = template.render(**context)

    pdf_path = Path(tmp_dir) / "relatorio.pdf"
    HTML(string=html_content).write_pdf(str(pdf_path))
    return pdf_path.read_bytes()


def _render_slides(context: dict[str, Any], tmp_dir: str) -> bytes:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Relatório da partida — SportsLyze"
    subtitle = f"{context['club_name']}"
    if context["opponent_name"]:
        subtitle += f" x {context['opponent_name']}"
    subtitle += f" — {context['match_date']} — Gerado em {context['generated_at']}"
    title_slide.placeholders[1].text = subtitle

    stats_slide = prs.slides.add_slide(prs.slide_layouts[5])
    stats_slide.shapes.title.text = "Estatísticas por jogador"

    players = context["players"]
    headers = ["Jogador", "Distância (km)", "Vel. média (km/h)", "Vel. máxima (km/h)"]
    rows = len(players) + 1
    table = stats_slide.shapes.add_table(
        rows, len(headers), Inches(0.5), Inches(1.6), Inches(9), Inches(0.4 * rows)
    ).table
    for col, header in enumerate(headers):
        table.cell(0, col).text = header
    for row_idx, player in enumerate(players, start=1):
        table.cell(row_idx, 0).text = player["name"]
        table.cell(row_idx, 1).text = f"{player['distance_km']:.2f}"
        table.cell(row_idx, 2).text = f"{player['avg_speed_kmh']:.1f}"
        table.cell(row_idx, 3).text = f"{player['max_speed_kmh']:.1f}"

    pptx_path = Path(tmp_dir) / "relatorio.pptx"
    prs.save(str(pptx_path))
    return pptx_path.read_bytes()


def _render_audio(context: dict[str, Any], tmp_dir: str) -> bytes:
    parts = [f"Relatório da partida do {context['club_name']}"]
    if context["opponent_name"]:
        parts.append(f"contra {context['opponent_name']}")
    if context["match_date"]:
        parts.append(f"em {context['match_date']}")
    text = ", ".join(parts) + ". "

    if context["players"]:
        text += "Estatísticas por jogador: "
        for player in context["players"]:
            text += (
                f"{player['name']}: {player['distance_km']:.1f} quilômetros percorridos, "
                f"velocidade média de {player['avg_speed_kmh']:.1f} quilômetros por hora. "
            )
    else:
        text += "Nenhuma estatística disponível ainda."

    audio_path = Path(tmp_dir) / "relatorio.mp3"
    gTTS(text=text, lang="pt-br").save(str(audio_path))
    return audio_path.read_bytes()
